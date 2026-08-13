import sys
from datetime import datetime
from pathlib import Path


def get_base_dir():
    if getattr(sys, "frozen", False):
        return Path(sys.executable).parent
    return Path(__file__).resolve().parent.parent


OUT_DIR = Path.home() / "Documents" / "Kaizumi"


def _ensure_out_dir() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    return OUT_DIR


def _sanitize(name: str) -> str:
    name = "".join(c for c in str(name) if c not in '\\/:*?"<>|').strip()
    return name or "document"


def create_presentation(title: str = "", slides=None, filename: str = None) -> str:
    """Create a .pptx from a title and a list of slides.

    slides: list of dicts {"title": str, "bullets": [str]} or plain strings.
    """
    from pptx import Presentation

    if not title:
        return "Please provide a presentation title, sir."
    slides = slides or []

    try:
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = title
        try:
            s.placeholders[1].text = datetime.now().strftime("%d.%m.%Y")
        except Exception:
            pass

        for idx, item in enumerate(slides, start=1):
            if isinstance(item, str):
                item = {"title": f"Slide {idx}", "bullets": [item]}
            slide_title = str(item.get("title") or f"Slide {idx}")
            bullets = item.get("bullets") or []
            if isinstance(bullets, str):
                bullets = [bullets]

            s = prs.slides.add_slide(content_layout)
            s.shapes.title.text = slide_title
            body = s.placeholders[1].text_frame
            first = True
            for b in bullets:
                p = body.paragraphs[0] if first else body.add_paragraph()
                first = False
                p.text = str(b)
                p.level = 0

        path = _ensure_out_dir() / f"{_sanitize(filename or title)}.pptx"
        prs.save(path)
        return f"Presentation created: {path}"
    except Exception as e:
        return f"Could not create presentation: {e}"


def create_spreadsheet(
    filename: str = "spreadsheet",
    headers=None,
    rows=None,
    sheet_name: str = "Sheet1",
) -> str:
    """Create an .xlsx with optional headers and data rows."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment
    from openpyxl.utils import get_column_letter

    headers = headers or []
    rows = rows or []

    try:
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name

        if headers:
            for c, h in enumerate(headers, start=1):
                cell = ws.cell(row=1, column=c, value=str(h))
                cell.font = Font(bold=True)
                cell.alignment = Alignment(horizontal="center")

        start_row = 2 if headers else 1
        for r, row in enumerate(rows, start=start_row):
            if isinstance(row, str):
                row = [row]
            elif isinstance(row, (int, float)):
                row = [row]
            else:
                row = list(row)
            for c, v in enumerate(row, start=1):
                ws.cell(row=r, column=c, value=v)

        max_cols = max(len(headers), *(len(r) if not isinstance(r, (str, int, float)) else 1 for r in rows), 1)
        for c in range(1, max_cols + 1):
            letter = get_column_letter(c)
            ws.column_dimensions[letter].width = 14

        path = _ensure_out_dir() / f"{_sanitize(filename)}.xlsx"
        wb.save(path)
        return f"Spreadsheet created: {path}"
    except Exception as e:
        return f"Could not create spreadsheet: {e}"