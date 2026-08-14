# actions/document_qa.py
# Kaizumi — question-answering over any document or web page
# (pdf, txt, md, html, docx, csv, pptx, xlsx, and http(s) URLs)

import os
import re
import sys
import zipfile
from pathlib import Path

from actions.pdf_reader import (
    _chunk_text,
    _pick_chunks,
    _ask_gemini,
)

try:
    from bs4 import BeautifulSoup
    _BS4 = True
except ImportError:
    _BS4 = False

try:
    import requests
    _REQUESTS = True
except ImportError:
    _REQUESTS = False


def get_base_dir() -> Path:
    if getattr(sys, "frozen", False):
        return Path(sys.executable).parent
    return Path(__file__).resolve().parent.parent


_URL_RE = re.compile(r"^(https?://)", re.IGNORECASE)

_SEARCH_DIRS = [
    Path.home() / "Documents",
    Path.home() / "Downloads",
    Path.home() / "Desktop",
    Path.home(),
]


def _resolve_source(raw: str) -> str:
    """Return either a local path (str(Path)) or a URL. None if not found."""
    raw = (raw or "").strip().strip('"')
    if not raw:
        return None
    if _URL_RE.match(raw):
        return raw

    cand = Path(raw)
    if cand.is_absolute() and cand.exists():
        return str(cand)
    if cand.exists():
        return str(cand)

    for d in _SEARCH_DIRS:
        for root, _dirs, files in os.walk(d):
            for f in files:
                if f.lower() == cand.name.lower():
                    return str(Path(root) / f)
            break
        p = d / cand.name
        if p.exists():
            return str(p)
    return None


def _extract_url(url: str) -> str:
    if not _REQUESTS:
        raise RuntimeError("requests is not installed.")
    headers = {"User-Agent": "Mozilla/5.0 (Kaizumi/assistant)"}
    resp = requests.get(url, headers=headers, timeout=20)
    resp.raise_for_status()
    if _BS4:
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        text = soup.get_text("\n", strip=True)
    else:
        text = re.sub(r"<[^>]+>", " ", resp.text)
    return re.sub(r"\n{3,}", "\n\n", text)


def _extract_docx(path: Path) -> str:
    with zipfile.ZipFile(str(path)) as z:
        xml = z.read("word/document.xml").decode("utf-8", errors="replace")
    xml = re.sub(r"</w:p>", "\n", xml)
    text = re.sub(r"<[^>]+>", "", xml)
    import html
    return html.unescape(text)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx is not installed.")
    prs = Presentation(str(path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("openpyxl is not installed.")
    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for ws in wb.worksheets:
        lines.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def _extract_text(source: str) -> str:
    if _URL_RE.match(source):
        return _extract_url(source)

    path = Path(source)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        from actions.pdf_reader import _extract_text as _pdf_text
        return _pdf_text(path)

    if suffix in (".docx",):
        return _extract_docx(path)

    if suffix in (".pptx", ".ppt"):
        return _extract_pptx(path)

    if suffix in (".xlsx", ".xls"):
        return _extract_xlsx(path)

    if suffix in (".html", ".htm"):
        raw = path.read_text(encoding="utf-8", errors="replace")
        if _BS4:
            soup = BeautifulSoup(raw, "html.parser")
            for tag in soup(["script", "style", "noscript"]):
                tag.decompose()
            return soup.get_text("\n", strip=True)
        return re.sub(r"<[^>]+>", " ", raw)

    # default: plain text (txt, md, csv, log, json, any)
    return path.read_text(encoding="utf-8", errors="replace")


def _describe(source: str) -> str:
    if _URL_RE.match(source):
        return source
    p = Path(source)
    return f"{p.name} ({p.stat().st_size // 1024} KB)"


def read_document(
    parameters: dict,
    response=None,
    player=None,
    session_memory=None,
) -> str:
    """Preview any local document or a web page."""
    params = parameters or {}
    source = _resolve_source(str(params.get("source", params.get("path", ""))))
    if not source:
        return ("I couldn't find that document, sir. Give me a file name, path, "
                "or a web address.")

    try:
        text = _extract_text(source)
    except Exception as e:
        return f"Could not read that source: {e}"
    if not text.strip():
        return f"'{_describe(source)}' has no extractable text, sir."

    preview = text[:1400] + ("…" if len(text) > 1400 else "")
    return f"Source: {_describe(source)} — {len(text)} characters.\n\n{preview}"


def document_qa(
    parameters: dict,
    response=None,
    player=None,
    session_memory=None,
) -> str:
    """Answer a question from any document or web page (lightweight RAG)."""
    params   = parameters or {}
    source   = _resolve_source(str(params.get("source", "")))
    question = str(params.get("question", params.get("query", ""))).strip()

    if not source:
        return ("I couldn't find that source, sir. Give a file name, path, or URL.")
    if not question:
        return "What would you like to know, sir?"

    try:
        text = _extract_text(source)
    except Exception as e:
        return f"Could not read that source: {e}"

    text_full = re.sub(r"\s+", " ", text).strip()
    if not text_full:
        return f"'{_describe(source)}' has no extractable text, sir."
    chunks = _chunk_text(text_full)
    picked = _pick_chunks(question, chunks)
    answer = _ask_gemini(question, "\n\n---\n\n".join(picked))
    return f"From {_describe(source)}:\n{answer}"